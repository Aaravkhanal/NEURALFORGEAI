from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Color Palette (NeuralForge Theme: Coral/Orange and Dark)
    theme_color = RGBColor(255, 68, 0) # #FF4400

    slides_content = [
        {
            "type": "title",
            "title": "NeuralForge AI",
            "subtitle": "Multi-LLM Agent System for Automated Machine Learning\n\nStreamlining Data to Deployment"
        },
        {
            "type": "bullet",
            "title": "1. What is NeuralForge AI?",
            "content": [
                "An intelligent, autonomous platform for end-to-end Machine Learning.",
                "Leverages Multiple Large Language Models (LLMs) to act as AI data scientists.",
                "Automates tedious tasks: data cleaning, feature engineering, and model selection.",
                "Bridges the gap between raw datasets and production-ready deployments."
            ]
        },
        {
            "type": "bullet",
            "title": "2. Core Platform Architecture",
            "content": [
                "Frontend: Modern, responsive Next.js application with Tailwind CSS & Shadcn UI.",
                "Backend: High-performance FastAPI server running Python.",
                "Database: Scalable architecture using SQLAlchemy (SQLite/PostgreSQL).",
                "AI Orchestration: LangChain & LangGraph for complex multi-agent workflows."
            ]
        },
        {
            "type": "bullet",
            "title": "3. Multi-LLM Agent Ecosystem",
            "content": [
                "Supports diverse LLM providers: OpenAI, Anthropic (Claude), Groq, and local models.",
                "Consultant Agent: Acts as your personal AI data science advisor.",
                "Data Cleaning Agent: Automatically detects and fixes anomalies in datasets.",
                "Research Agent: Fetches the latest ML papers and state-of-the-art techniques."
            ]
        },
        {
            "type": "bullet",
            "title": "4. Intelligent Data Processing",
            "content": [
                "Automated Missing Value Imputation: AI decides the best strategy (mean, median, advanced).",
                "Outlier Detection & Handling: Ensures robust datasets for training.",
                "Feature Engineering: Automatically discovers and creates predictive features.",
                "Data Profiling: Generates comprehensive statistical reports before training."
            ]
        },
        {
            "type": "bullet",
            "title": "5. Advanced AutoML & Training",
            "content": [
                "Integrated with top-tier frameworks: XGBoost, LightGBM, CatBoost, and Scikit-Learn.",
                "Hyperparameter Optimization: Automatically tunes models for maximum accuracy.",
                "Continuous Monitoring: Tracks training jobs and logs performance metrics in real-time.",
                "Model Registry: Version control for your trained ML models."
            ]
        },
        {
            "type": "bullet",
            "title": "6. Probabilistic Approaches & Techniques",
            "content": [
                "Bayesian Hyperparameter Optimization: Probabilistic search for optimal model tuning.",
                "Ensemble Probabilities: Tree-based models (XGBoost, LightGBM) output calibrated probability scores.",
                "Synthetic Data Generation: Generative probabilistic sampling to augment imbalanced datasets.",
                "LLM Autoregression: Core AI agents rely on probabilistic token generation for decision making."
            ]
        },
        {
            "type": "bullet",
            "title": "7. Model Transparency & Explainability",
            "content": [
                "No 'Black Box' Models: Full visibility into how decisions are made.",
                "SHAP Integration: Analyzes global and local feature importance.",
                "LIME Integration: Explains individual predictions for deep insights.",
                "Interactive Visualizations: Beautiful charts to interpret model behavior."
            ]
        },
        {
            "type": "bullet",
            "title": "8. Interactive AI Playground",
            "content": [
                "Test Models Instantly: An interactive sandbox to run predictions on your trained models.",
                "Dataset Chat: Talk directly to your CSV/Excel files using LLMs.",
                "Real-time Feedback: Instantly see how changes in inputs affect the prediction.",
                "Compare Models: Head-to-head comparison of different algorithms."
            ]
        },
        {
            "type": "bullet",
            "title": "9. Seamless Deployment & Export",
            "content": [
                "One-Click Deployment: Move models from the sandbox to production effortlessly.",
                "Code Generation: NeuralForge writes the API code (Python/FastAPI) to host your model.",
                "Export Options: Download cleaned datasets, trained weights, or containerized apps.",
                "API Ready: Automatically exposes REST endpoints for integration."
            ]
        },
        {
            "type": "bullet",
            "title": "10. Security & Enterprise Readiness",
            "content": [
                "Secure Authentication: JWT-based auth and secure credential management.",
                "Isolated Environments: Projects are siloed to ensure data privacy.",
                "Scalable Task Queues: Celery & Redis integration for heavy, long-running ML jobs.",
                "Comprehensive Logging: Detailed logs for compliance and debugging."
            ]
        },
        {
            "type": "title",
            "title": "11. Conclusion & Future Roadmap",
            "subtitle": "NeuralForge AI empowers anyone to build state-of-the-art ML systems.\n\nThank you!\n\nQuestions?"
        }
    ]

    for slide_data in slides_content:
        if slide_data["type"] == "title":
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.color.rgb = theme_color
            title.text_frame.paragraphs[0].font.bold = True
            
            subtitle.text = slide_data["subtitle"]
            
        elif slide_data["type"] == "bullet":
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = slide_data["title"]
            title_shape.text_frame.paragraphs[0].font.color.rgb = theme_color
            
            tf = body_shape.text_frame
            tf.text = slide_data["content"][0]
            
            for point in slide_data["content"][1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    prs.save('NeuralForge_AI_Presentation.pptx')
    print("Presentation saved as NeuralForge_AI_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
